from abc import ABC, abstractmethod
import json

class DocumentProcessor(ABC):
    @abstractmethod
    def read_document(self, input_path, output_path) -> list:
        pass

    @abstractmethod
    def modify_document(self, input_path, output_path, replacements):
        pass


class TxtDocumentProcessor(DocumentProcessor):
    def read_document(self, input_path, output_path) -> list:
        segments = []
        with open(input_path, 'r', encoding='utf-8') as infile:
            segment_id = 0
            for line in infile:
                if line.strip() == "":
                    continue
                segments.append({
                    "id": segment_id,
                    "original_text": line.strip()
                })
                segment_id += 1
        # with open(output_path, "w", encoding='utf-8') as outfile:
        #     json.dump(segments, outfile, ensure_ascii=False, indent=4)
        return segments

    def modify_document(self, input_path, output_path, replacements):
        lines = []
        with open(input_path, 'r', encoding='utf-8') as infile:
            segment_id = 0
            for line in infile:
                if line.strip() == "":
                    lines.append(line)
                    continue
                modified_line = replacements[segment_id].get('translated_text', line.strip())
                lines.append(modified_line + "\n")
                segment_id += 1
        with open(output_path, 'w', encoding='utf-8') as outfile:
            outfile.writelines(lines)


class DocxDocumentProcessor(DocumentProcessor):
    def read_document(self, input_path, output_path) -> list:
        from docx import Document
        doc = Document(input_path)
        segments = []
        segment_id = 0
        for paragraph in doc.paragraphs:
            if paragraph.text.strip() == "":
                continue
            segments.append({
                "id": segment_id,
                "original_text": paragraph.text.strip()
            })
            segment_id += 1
        json.dump(segments, open(output_path, "w", encoding='utf-8'), ensure_ascii=False, indent=4)
        return segments

    def modify_document(self, input_path, output_path, replacements):
        from docx import Document
        doc = Document(input_path)
        segment_id = 0
        for paragraph in doc.paragraphs:
            if paragraph.text.strip() == "":
                continue
            paragraph.text = replacements[segment_id]['translated_text']
            segment_id += 1
        doc.save(output_path)


class PptxDocumentProcessor(DocumentProcessor):
    def read_document(self, input_path, output_path) -> list:
        from pptx import Presentation
        prs = Presentation(input_path)
        segment_id = 0
        segments = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                original_text = text_frame.text
                if original_text.strip() == "":
                    continue
                segments.append({
                    "id": segment_id,
                    "original_text": original_text.strip()
                })
                segment_id += 1
        json.dump(segments, open(output_path, "w", encoding='utf-8'), ensure_ascii=False, indent=4)
        return segments

    def modify_document(self, input_path, output_path, replacements):
        from pptx import Presentation
        prs = Presentation(input_path)
        segment_id = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                if text_frame.text.strip() == "":
                    continue
                text_frame.text = replacements[segment_id]['translated_text']
                segment_id += 1
        prs.save(output_path)


class PdfDocumentProcessor(DocumentProcessor):
    def read_document(self, input_path, output_path) -> list:
        import fitz  # PyMuPDF
        doc = fitz.open(input_path)
        segment_id = 0
        segments = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            blocks = page.get_text("blocks")
            for block in blocks:
                text = block[4]
                if text.strip() == "":
                    continue
                segments.append({
                    "id": segment_id,
                    "original_text": text.strip()
                })
                segment_id += 1
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(segments, f, ensure_ascii=False, indent=4)
        return segments

    def modify_document(self, input_path, output_path, replacements):
        lines = ''
        for item in replacements:
            lines += item['translated_text'] + '\n'
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(lines)


# gRPC 服务实现
from concurrent import futures
import grpc
import document_processor_pb2
import document_processor_pb2_grpc

class DocumentProcessorServicer(document_processor_pb2_grpc.DocumentProcessorServicer):
    def __init__(self):
        self.processors = {
            'txt': TxtDocumentProcessor(),
            'docx': DocxDocumentProcessor(),
            'pptx': PptxDocumentProcessor(),
            'pdf': PdfDocumentProcessor()
        }

    def ProcessDocument(self, request, context):
        processor = self.processors.get(request.file_type)
        if not processor:
            context.set_code(grpc.StatusCode.INVALID_ARGUMENT)
            context.set_details('Unsupported file type')
            return document_processor_pb2.ProcessDocumentResponse()

        segments = processor.read_document(request.input_path, request.output_path)
        return document_processor_pb2.ProcessDocumentResponse(segments=json.dumps(segments))

    def ModifyDocument(self, request, context):
        processor = self.processors.get(request.file_type)
        if not processor:
            context.set_code(grpc.StatusCode.INVALID_ARGUMENT)
            context.set_details('Unsupported file type')
            return document_processor_pb2.ModifyDocumentResponse()

        replacements = json.loads(request.replacements)
        processor.modify_document(request.input_path, request.output_path, replacements)
        return document_processor_pb2.ModifyDocumentResponse(success=True)


def serve():
    server = grpc.server(futures.ThreadPoolExecutor(max_workers=10))
    document_processor_pb2_grpc.add_DocumentProcessorServicer_to_server(DocumentProcessorServicer(), server)
    server.add_insecure_port('[::]:50051')
    server.start()
    server.wait_for_termination()


if __name__ == '__main__':
    serve()
