from abc import ABC, abstractmethod
import json
from flask import Flask, request, jsonify, send_file
from io import BytesIO
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

app = Flask(__name__)


class DocumentProcessor(ABC):
    @abstractmethod
    def read_document(self, content: str) -> list:
        pass

    @abstractmethod
    def modify_document(self, content: str, replacements: list) -> str:
        pass


class TxtDocumentProcessor(DocumentProcessor):
    def read_document(self, content: BytesIO) -> list:
        content = content.getvalue().decode('utf-8')

        segments = []
        segment_id = 0
        for line in content.splitlines():
            if line.strip() == "":
                continue
            segments.append({
                "relative_id": segment_id,
                "source_text": line.strip()
            })
            segment_id += 1
        return segments

    def modify_document(self, content: BytesIO, replacements: list) -> BytesIO:
        sorted_content = sorted(replacements, key=lambda x: x['relative_id'])
        lines = [item['target_text'] for item in sorted_content]
        modified_content = "\n".join(lines)
        return BytesIO(modified_content.encode('utf-8'))


class DocxDocumentProcessor(DocumentProcessor):
    def read_document(self, content: BytesIO) -> list:
        doc = Document(content)
        segments = []
        segment_id = 0
        for paragraph in doc.paragraphs:
            if paragraph.text.strip() == "":
                continue
            segments.append({
                "relative_id": segment_id,
                "source_text": paragraph.text.strip()
            })
            segment_id += 1
        return segments

    def modify_document(self, content: BytesIO, replacements: list) -> BytesIO:
        doc = Document(content)
        segment_id = 0
        for paragraph in doc.paragraphs:
            if paragraph.text.strip() == "":
                continue
            paragraph.text = replacements[segment_id]['target_text']
            segment_id += 1
        output = BytesIO()
        doc.save(output)
        output.seek(0)  # 将指针移回开头以便后续读取

        return output


class PptxDocumentProcessor(DocumentProcessor):
    def read_document(self, content: BytesIO) -> list:
        prs = Presentation(content)
        segments = []
        segment_id = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                source_text = text_frame.text
                if source_text.strip() == "":
                    continue
                segments.append({
                    "relative_id": segment_id,
                    "source_text": source_text.strip()
                })
                segment_id += 1
        return segments

    def modify_document(self, content: BytesIO, replacements: list) -> BytesIO:
        prs = Presentation(content)
        segment_id = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                if text_frame.text.strip() == "":
                    continue
                text_frame.text = replacements[segment_id]['target_text']
                segment_id += 1
        output = BytesIO()
        prs.save(output)
        output.seek(0)  # 将指针移回开头以便后续读取
        prs.save('./uu.pptx')
        return output


class PdfDocumentProcessor(DocumentProcessor):
    def read_document(self, content: str) -> list:
        doc = fitz.open("pdf", content)
        segments = []
        segment_id = 0
        for page_num in range(len(doc)):
            page = doc[page_num]
            blocks = page.get_text("blocks")
            for block in blocks:
                text = block[4]
                if text.strip() == "":
                    continue
                segments.append({
                    "relative_id": segment_id,
                    "source_text": text.strip()
                })
                segment_id += 1
        return segments

    def modify_document(self, content: str, replacements: list) -> BytesIO:
        lines = '\n'.join([item['translated_text'] for item in replacements])
        return BytesIO(lines.encode('utf-8'))


processors = {
    'txt': TxtDocumentProcessor(),
    'docx': DocxDocumentProcessor(),
    'pptx': PptxDocumentProcessor(),
    'pdf': PdfDocumentProcessor()
}


@app.route('/process_document', methods=['POST'])
def process_document():
    file = request.files['file']
    file_type = file.filename.split('.')[-1]

    processor = processors.get(file_type)
    if not processor:
        return jsonify({'error': 'Unsupported file type'}), 400

    # file.save('./u.pptx')
    content = BytesIO(file.read())
    segments = processor.read_document(content)
    return jsonify({'segments': segments})


@app.route('/modify_document', methods=['POST'])
def modify_document():
    file = request.files['file']
    file_type = file.filename.split('.')[-1]
    replacements = json.loads(request.form['replacements'])

    processor = processors.get(file_type)
    if not processor:
        return jsonify({'error': 'Unsupported file type'}), 400

    content = BytesIO(file.read())
    modified_content = processor.modify_document(content, replacements)
    if file_type == 'pdf':
        response = send_file(modified_content, as_attachment=True, download_name=f'modified_{file.filename}.txt')
    else:
        response = send_file(modified_content, as_attachment=True, download_name=f'modified_{file.filename}')
    return response


if __name__ == '__main__':
    # 打印当前运行目录
    import os
    current_directory = os.getcwd()
    print(f'Current directory: {current_directory}')

    # 打印当前文件夹下的文件（递归）
    print('Files in current directory (recursive):')
    for root, dirs, files in os.walk(current_directory):
        for file in files:
            print(os.path.join(root, file))

    app.run(host='0.0.0.0', port=5000)
